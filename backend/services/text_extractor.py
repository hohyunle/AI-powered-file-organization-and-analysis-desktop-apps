"""
텍스트 추출 모듈

다양한 파일 형식에서 텍스트를 추출합니다.
- PDF: PyMuPDF 사용
- PowerPoint: python-pptx 사용  
- Word: python-docx 사용
- Excel: openpyxl 사용
- 이미지: pytesseract (OCR) 사용
"""

import asyncio
from pathlib import Path
from typing import Dict, Any, Optional
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
from PIL import Image
import pytesseract
from loguru import logger


class TextExtractor:
    """텍스트 추출 클래스"""
    
    def __init__(self):
        # Tesseract 경로 설정 (Windows)
        try:
            pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        except:
            logger.warning("Tesseract OCR 경로를 찾을 수 없습니다. OCR 기능이 제한될 수 있습니다.")
    
    async def extract_text(self, file_path: str) -> Dict[str, Any]:
        """파일에서 텍스트 추출"""
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                return {"error": "파일이 존재하지 않습니다."}
            
            file_extension = file_path.suffix.lower()
            
            # 파일 형식별 추출
            if file_extension == '.pdf':
                return await self._extract_pdf_text(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                return await self._extract_pptx_text(file_path)
            elif file_extension in ['.docx', '.doc']:
                return await self._extract_docx_text(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                return await self._extract_xlsx_text(file_path)
            elif file_extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
                return await self._extract_image_text(file_path)
            elif file_extension == '.txt':
                return await self._extract_txt_text(file_path)
            else:
                return {"error": f"지원하지 않는 파일 형식: {file_extension}"}
                
        except Exception as e:
            logger.error(f"텍스트 추출 오류 ({file_path}): {e}")
            return {"error": str(e)}
    
    async def _extract_pdf_text(self, file_path: Path) -> Dict[str, Any]:
        """PDF 파일에서 텍스트 추출"""
        try:
            doc = fitz.open(str(file_path))
            text_content = []
            page_count = len(doc)
            
            for page_num in range(page_count):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    text_content.append(text.strip())
            
            doc.close()
            
            # 전체 텍스트 결합
            full_text = "\n\n".join(text_content)
            
            return {
                "text_snippet": full_text[:1000],  # 처음 1000자만 저장
                "full_text": full_text,
                "page_count": page_count,
                "ocr_used": False,
                "lang": "ko",  # 기본값, 향후 언어 감지 추가
                "meta_json": {
                    "page_count": page_count,
                    "file_type": "PDF",
                    "extraction_method": "PyMuPDF"
                }
            }
            
        except Exception as e:
            logger.error(f"PDF 추출 오류: {e}")
            return {"error": f"PDF 추출 실패: {e}"}
    
    async def _extract_pptx_text(self, file_path: Path) -> Dict[str, Any]:
        """PowerPoint 파일에서 텍스트 추출"""
        try:
            prs = Presentation(str(file_path))
            text_content = []
            slide_count = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_content.append(f"슬라이드 {slide_num + 1}: " + " ".join(slide_text))
            
            full_text = "\n\n".join(text_content)
            
            return {
                "text_snippet": full_text[:1000],
                "full_text": full_text,
                "slide_count": slide_count,
                "ocr_used": False,
                "lang": "ko",
                "meta_json": {
                    "slide_count": slide_count,
                    "file_type": "PowerPoint",
                    "extraction_method": "python-pptx"
                }
            }
            
        except Exception as e:
            logger.error(f"PowerPoint 추출 오류: {e}")
            return {"error": f"PowerPoint 추출 실패: {e}"}
    
    async def _extract_docx_text(self, file_path: Path) -> Dict[str, Any]:
        """Word 문서에서 텍스트 추출"""
        try:
            doc = Document(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n".join(paragraphs)
            
            return {
                "text_snippet": full_text[:1000],
                "full_text": full_text,
                "paragraph_count": len(paragraphs),
                "ocr_used": False,
                "lang": "ko",
                "meta_json": {
                    "paragraph_count": len(paragraphs),
                    "file_type": "Word",
                    "extraction_method": "python-docx"
                }
            }
            
        except Exception as e:
            logger.error(f"Word 추출 오류: {e}")
            return {"error": f"Word 추출 실패: {e}"}
    
    async def _extract_xlsx_text(self, file_path: Path) -> Dict[str, Any]:
        """Excel 파일에서 텍스트 추출"""
        try:
            workbook = load_workbook(str(file_path))
            text_content = []
            sheet_count = len(workbook.sheetnames)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []
                
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value:
                            row_text.append(str(cell.value))
                    if row_text:
                        sheet_text.append(" ".join(row_text))
                
                if sheet_text:
                    text_content.append(f"시트 '{sheet_name}': " + " ".join(sheet_text))
            
            full_text = "\n\n".join(text_content)
            
            return {
                "text_snippet": full_text[:1000],
                "full_text": full_text,
                "sheet_count": sheet_count,
                "ocr_used": False,
                "lang": "ko",
                "meta_json": {
                    "sheet_count": sheet_count,
                    "sheet_names": workbook.sheetnames,
                    "file_type": "Excel",
                    "extraction_method": "openpyxl"
                }
            }
            
        except Exception as e:
            logger.error(f"Excel 추출 오류: {e}")
            return {"error": f"Excel 추출 실패: {e}"}
    
    async def _extract_image_text(self, file_path: Path) -> Dict[str, Any]:
        """이미지에서 OCR로 텍스트 추출"""
        try:
            # 이미지 열기
            image = Image.open(str(file_path))
            
            # OCR 실행
            text = pytesseract.image_to_string(image, lang='kor+eng')
            
            return {
                "text_snippet": text[:1000],
                "full_text": text,
                "ocr_used": True,
                "lang": "ko",
                "meta_json": {
                    "image_size": image.size,
                    "image_mode": image.mode,
                    "file_type": "Image",
                    "extraction_method": "Tesseract OCR"
                }
            }
            
        except Exception as e:
            logger.error(f"이미지 OCR 오류: {e}")
            return {"error": f"이미지 OCR 실패: {e}"}
    
    async def _extract_txt_text(self, file_path: Path) -> Dict[str, Any]:
        """텍스트 파일에서 내용 추출"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            return {
                "text_snippet": text[:1000],
                "full_text": text,
                "ocr_used": False,
                "lang": "ko",
                "meta_json": {
                    "file_type": "Text",
                    "extraction_method": "direct_read"
                }
            }
            
        except Exception as e:
            logger.error(f"텍스트 파일 읽기 오류: {e}")
            return {"error": f"텍스트 파일 읽기 실패: {e}"}
